"""Generate project showcase PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x3E)   # slide background / headers
BLUE      = RGBColor(0x1A, 0x56, 0xDB)   # accent
LIGHT_BG  = RGBColor(0xF4, 0xF7, 0xFF)   # light panel bg
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x64, 0x74, 0x8B)
GOLD      = RGBColor(0xF5, 0xA6, 0x23)   # highlight numbers

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helper utilities ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_w=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0, bullet=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    if bullet:
        p.text = text
    else:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    if not bullet:
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return p


def navy_slide(slide):
    """Fill slide background with navy."""
    add_rect(slide, 0, 0, W, H, fill=NAVY)


def slide_title_bar(slide, title, subtitle=None):
    """Top bar with title."""
    add_rect(slide, 0, 0, W, Inches(1.1), fill=BLUE)
    add_text(slide, title,
             Inches(0.4), Inches(0.12), Inches(12), Inches(0.7),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.75), Inches(12), Inches(0.4),
                 size=14, color=RGBColor(0xBF, 0xD9, 0xFF))


def footer(slide, text="EMS Competitive Intelligence Platform  ·  Flex Ltd  ·  2026"):
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), fill=RGBColor(0x08, 0x10, 0x28))
    add_text(slide, text,
             Inches(0.3), Inches(7.12), Inches(12), Inches(0.35),
             size=10, color=GRAY, align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
navy_slide(s1)

# Decorative accent bar
add_rect(s1, 0, Inches(2.8), Inches(0.12), Inches(2.0), fill=GOLD)

# Main title
add_text(s1, "EMS Competitive Intelligence Platform",
         Inches(0.5), Inches(2.2), Inches(11), Inches(1.2),
         size=40, bold=True, color=WHITE)

# Subtitle
add_text(s1, "An AI-Powered Research System for EMS Sector Analysis",
         Inches(0.5), Inches(3.35), Inches(11), Inches(0.6),
         size=22, color=RGBColor(0xBF, 0xD9, 0xFF))

# Meta line
add_text(s1, "Project Showcase  ·  Sponsoring Company: Flex Ltd  ·  May 2026",
         Inches(0.5), Inches(4.1), Inches(11), Inches(0.4),
         size=14, color=GRAY)

# Bottom accent
add_rect(s1, 0, Inches(7.1), W, Inches(0.4), fill=BLUE)
add_text(s1, "10-Minute Presentation  |  Includes 3-Min Platform Demo",
         Inches(0.4), Inches(7.12), Inches(12), Inches(0.35),
         size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Project Scope
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
navy_slide(s2)
slide_title_bar(s2, "Project Scope", "What problem we solved — and what we built")
footer(s2)

# Left panel — The Challenge
add_rect(s2, Inches(0.3), Inches(1.3), Inches(5.9), Inches(5.5),
         fill=RGBColor(0x13, 0x26, 0x52))

add_text(s2, "THE CHALLENGE",
         Inches(0.55), Inches(1.45), Inches(5.4), Inches(0.4),
         size=12, bold=True, color=GOLD)

challenges = [
    "Flex operates in a $200B+ EMS market with 5 major competitors",
    "No centralized system to monitor competitor strategy or CapEx trends in real time",
    "Analysts spend hours manually reading SEC filings, earnings calls, and news",
    "Hyperscaler AI investment signals — critical for EMS demand forecasting — scattered across sources",
]
y = Inches(1.95)
for c in challenges:
    add_rect(s2, Inches(0.55), y, Inches(0.06), Inches(0.06), fill=GOLD)
    add_text(s2, c, Inches(0.72), y - Inches(0.05), Inches(5.3), Inches(0.55),
             size=13, color=WHITE)
    y += Inches(0.72)

# Right panel — The Solution
add_rect(s2, Inches(6.5), Inches(1.3), Inches(6.55), Inches(5.5),
         fill=RGBColor(0x13, 0x26, 0x52))

add_text(s2, "OUR SOLUTION",
         Inches(6.75), Inches(1.45), Inches(6.0), Inches(0.4),
         size=12, bold=True, color=GOLD)

add_text(s2, "A unified AI-powered competitive intelligence platform",
         Inches(6.75), Inches(1.9), Inches(5.9), Inches(0.5),
         size=14, bold=True, color=WHITE)

stats = [
    ("11", "Companies tracked\n(6 EMS + 5 Hyperscalers)"),
    ("8",  "Functional modules"),
    ("20+", "Years of financial history\nvia SEC EDGAR"),
]
sx = Inches(6.75)
for num, label in stats:
    add_text(s2, num, sx, Inches(2.6), Inches(1.5), Inches(0.8),
             size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s2, label, sx, Inches(3.2), Inches(1.6), Inches(0.7),
             size=11, color=RGBColor(0xBF, 0xD9, 0xFF), align=PP_ALIGN.CENTER)
    sx += Inches(1.9)

add_text(s2, "Data sources: SEC EDGAR · Real-time News · yfinance · Earnings Transcripts",
         Inches(6.75), Inches(4.15), Inches(5.9), Inches(0.4),
         size=12, color=GRAY)

modules = ["📰 News Intelligence", "🧠 Analyst View", "💬 AI Chat", "🏢 Companies Hub",
           "📈 AI Investments", "🗺️ Facilities Map", "📅 Calendar", "🗄️ Data Center"]
mx, my = Inches(6.75), Inches(4.65)
for i, m in enumerate(modules):
    col = i % 2
    row = i // 2
    add_text(s2, m,
             mx + col * Inches(2.85), my + row * Inches(0.42),
             Inches(2.7), Inches(0.38),
             size=12, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Work Steps: Data Infrastructure
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
navy_slide(s3)
slide_title_bar(s3, "Work Steps  —  Part 1: Data Foundation",
                "How we built the intelligence backbone")
footer(s3)

steps = [
    {
        "num": "01",
        "title": "Data Ingestion",
        "points": [
            "SEC EDGAR Company Facts API → 20+ years financial history",
            "News APIs → real-time competitor coverage across 11 companies",
            "yfinance → market data & fallback financials",
            "Earnings call transcripts → automated ingestion pipeline",
        ],
    },
    {
        "num": "02",
        "title": "Data Normalization",
        "points": [
            "6 EMS companies × 6 different fiscal year endings → unified FY/Q labels",
            "XBRL concept mapping: GAAP field names → standardized schema",
            "Cross-company CapEx comparison with apples-to-apples alignment",
            "ChromaDB vector store: SEC filings chunked for semantic search",
        ],
    },
    {
        "num": "03",
        "title": "Storage Layer",
        "points": [
            "SQLite financial cache → sub-second queries, no LLM needed for numbers",
            "ChromaDB → document retrieval for qualitative filing analysis",
            "JSON caches → news, analytics results, preset questions",
        ],
    },
]

panel_w = Inches(3.9)
px = Inches(0.3)
for step in steps:
    add_rect(s3, px, Inches(1.3), panel_w, Inches(5.55),
             fill=RGBColor(0x13, 0x26, 0x52))
    # Number badge
    add_rect(s3, px + Inches(0.15), Inches(1.45), Inches(0.55), Inches(0.55), fill=BLUE)
    add_text(s3, step["num"],
             px + Inches(0.15), Inches(1.43), Inches(0.55), Inches(0.55),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s3, step["title"],
             px + Inches(0.82), Inches(1.5), Inches(2.9), Inches(0.45),
             size=17, bold=True, color=WHITE)
    py = Inches(2.15)
    for pt in step["points"]:
        add_rect(s3, px + Inches(0.25), py + Inches(0.12),
                 Inches(0.07), Inches(0.07), fill=GOLD)
        add_text(s3, pt,
                 px + Inches(0.42), py, Inches(3.35), Inches(0.65),
                 size=12, color=RGBColor(0xD0, 0xDC, 0xFF))
        py += Inches(0.75)
    px += panel_w + Inches(0.22)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Work Steps: AI Intelligence Layer
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
navy_slide(s4)
slide_title_bar(s4, "Work Steps  —  Part 2: AI & Analytics Layer",
                "From raw data to actionable intelligence")
footer(s4)

modules_ai = [
    {
        "icon": "🔍",
        "title": "RAG Pipeline",
        "desc": "Retrieval-Augmented Generation over SEC filings. Query → embedding → ChromaDB → LLM synthesis. Sources always cited; grounded in actual documents.",
    },
    {
        "icon": "⚡",
        "title": "Financial Cache Short-circuit",
        "desc": "Numeric queries (CapEx, Revenue…) answered directly from SQLite. Bypasses LLM entirely → <1 second response. Covers 11 companies × 3 statements × 20+ years.",
    },
    {
        "icon": "📊",
        "title": "Analytics Engine",
        "desc": "Sentiment analysis on earnings transcripts. Statistical anomaly detection for CapEx outliers. AI strategy classifier. Geographic footprint scoring.",
    },
    {
        "icon": "🖥️",
        "title": "Next.js Dashboard",
        "desc": "8 functional modules. Real-time streaming AI responses. Persistent session management. Structured table rendering. Dark / light mode.",
    },
]

bw = Inches(6.1)
positions = [(Inches(0.25), Inches(1.3)), (Inches(6.7), Inches(1.3)),
             (Inches(0.25), Inches(4.2)), (Inches(6.7), Inches(4.2))]

for (bx, by), m in zip(positions, modules_ai):
    add_rect(s4, bx, by, bw, Inches(2.65), fill=RGBColor(0x13, 0x26, 0x52))
    add_text(s4, m["icon"] + "  " + m["title"],
             bx + Inches(0.2), by + Inches(0.15), bw - Inches(0.3), Inches(0.5),
             size=18, bold=True, color=WHITE)
    add_rect(s4, bx + Inches(0.2), by + Inches(0.68), bw - Inches(0.4),
             Inches(0.03), fill=BLUE)
    add_text(s4, m["desc"],
             bx + Inches(0.2), by + Inches(0.8), bw - Inches(0.4), Inches(1.7),
             size=13, color=RGBColor(0xBF, 0xD9, 0xFF))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Key Deliverables (transition to video)
# ─────────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
navy_slide(s5)
slide_title_bar(s5, "Key Deliverables", "8 modules — see them in action")
footer(s5)

deliverables = [
    ("📰", "News Intelligence Feed",     "Real-time competitor news, topic filtering"),
    ("🧠", "AI Analyst View",            "Weekly themes, sentiment, executive summaries"),
    ("💬", "AI Chat",                    "Natural language financial queries, table output"),
    ("🏢", "Companies Hub",              "Per-company deep dives: financials, CapEx, hiring"),
    ("📈", "Hyperscaler AI Investments", "Big Tech CapEx trends & YoY growth tracking"),
    ("🗺️", "Global Facilities Map",     "Interactive globe of EMS manufacturing locations"),
    ("📅", "Earnings & Events Calendar", "Upcoming earnings, analyst events, key dates"),
    ("🗄️", "Data Center",               "Full SEC filing library, downloadable documents"),
]

col_w = Inches(6.1)
for i, (icon, title, desc) in enumerate(deliverables):
    col = i % 2
    row = i // 2
    bx = Inches(0.25) + col * Inches(6.55)
    by = Inches(1.35) + row * Inches(1.42)
    add_rect(s5, bx, by, col_w, Inches(1.28), fill=RGBColor(0x13, 0x26, 0x52))
    add_rect(s5, bx, by, Inches(0.08), Inches(1.28), fill=BLUE)
    add_text(s5, icon + "  " + title,
             bx + Inches(0.2), by + Inches(0.08), col_w - Inches(0.3), Inches(0.45),
             size=15, bold=True, color=WHITE)
    add_text(s5, desc,
             bx + Inches(0.2), by + Inches(0.55), col_w - Inches(0.3), Inches(0.55),
             size=12, color=GRAY)

# Video cue
add_rect(s5, Inches(3.5), Inches(6.9), Inches(6.3), Inches(0.45), fill=BLUE)
add_text(s5, "▶  Next: 3-minute platform demo",
         Inches(3.5), Inches(6.9), Inches(6.3), Inches(0.45),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — [VIDEO PLACEHOLDER]
# ─────────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, W, H, fill=RGBColor(0x05, 0x0A, 0x1A))
add_text(s6, "[ PLATFORM DEMO VIDEO ]",
         Inches(2), Inches(3.2), Inches(9.33), Inches(0.8),
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s6, "Insert screen_recording.webm / final_output.mp4 here  —  ~3 minutes",
         Inches(2), Inches(4.0), Inches(9.33), Inches(0.5),
         size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_rect(s6, Inches(1.5), Inches(2.5), Inches(10.33), Inches(3.5),
         line_color=RGBColor(0x1A, 0x56, 0xDB), line_w=Pt(2))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Project Benefits to Flex
# ─────────────────────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
navy_slide(s7)
slide_title_bar(s7, "Project Benefits to Flex", "Efficiency · Intelligence Quality · Strategic Advantage")
footer(s7)

benefits = [
    {
        "icon": "⏱",
        "title": "Efficiency",
        "color": BLUE,
        "points": [
            "Financial data retrieval: hours of manual reading  →  <1 second",
            "Competitor monitoring: manual daily scan  →  automated, real-time",
            "Earnings call analysis: 60-page transcript  →  AI summary in seconds",
        ],
    },
    {
        "icon": "📊",
        "title": "Intelligence Quality",
        "color": RGBColor(0x05, 0x96, 0x69),
        "points": [
            "20+ years of CapEx history vs. 5 years from standard data tools",
            "Cross-company fiscal year normalization — apples-to-apples comparison",
            "Statistical anomaly detection flags unusual strategic shifts early",
        ],
    },
    {
        "icon": "🔭",
        "title": "Strategic Advantage",
        "color": GOLD,
        "points": [
            "Hyperscaler CapEx trends directly inform EMS demand forecasting",
            "Real-time signals from 11 companies in one unified view",
            "Platform is extensible — new companies & data sources add easily",
        ],
    },
]

bw = Inches(4.0)
bx = Inches(0.25)
for b in benefits:
    add_rect(s7, bx, Inches(1.3), bw, Inches(5.5), fill=RGBColor(0x13, 0x26, 0x52))
    add_rect(s7, bx, Inches(1.3), bw, Inches(0.55), fill=b["color"])
    add_text(s7, b["icon"] + "  " + b["title"],
             bx + Inches(0.15), Inches(1.32), bw - Inches(0.2), Inches(0.5),
             size=18, bold=True, color=WHITE)
    py = Inches(2.05)
    for pt in b["points"]:
        add_rect(s7, bx + Inches(0.2), py + Inches(0.14),
                 Inches(0.08), Inches(0.08), fill=b["color"])
        add_text(s7, pt,
                 bx + Inches(0.38), py, bw - Inches(0.55), Inches(0.75),
                 size=13, color=RGBColor(0xD0, 0xDC, 0xFF))
        py += Inches(0.85)
    bx += bw + Inches(0.22)

# Callout quote
add_rect(s7, Inches(0.25), Inches(6.55), Inches(12.85), Inches(0.65),
         fill=RGBColor(0x08, 0x20, 0x55))
add_text(s7,
         "\"This platform gives Flex's research team institutional-grade competitive intelligence — built entirely on public data.\"",
         Inches(0.5), Inches(6.58), Inches(12.35), Inches(0.6),
         size=13, italic=True, color=RGBColor(0xBF, 0xD9, 0xFF),
         align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Summary & Q&A
# ─────────────────────────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
navy_slide(s8)
add_rect(s8, 0, 0, W, Inches(1.1), fill=BLUE)
add_text(s8, "Thank You",
         Inches(0.4), Inches(0.12), Inches(12), Inches(0.85),
         size=34, bold=True, color=WHITE)
footer(s8)

add_text(s8, "What we built — in three lines:",
         Inches(0.5), Inches(1.4), Inches(12), Inches(0.45),
         size=16, color=GOLD, bold=True)

summary_pts = [
    "Built an end-to-end AI competitive intelligence platform in one quarter",
    "11 companies  ·  8 modules  ·  20+ years of data  ·  <1 sec query response",
    "Turns scattered public data into real-time, actionable insight for Flex",
]
sy = Inches(1.95)
for pt in summary_pts:
    add_rect(s8, Inches(0.5), sy + Inches(0.12), Inches(0.1), Inches(0.1), fill=GOLD)
    add_text(s8, pt, Inches(0.72), sy, Inches(11.5), Inches(0.5),
             size=18, color=WHITE)
    sy += Inches(0.62)

add_rect(s8, Inches(3.5), Inches(4.2), Inches(6.3), Inches(2.2),
         fill=RGBColor(0x13, 0x26, 0x52))
add_text(s8, "Questions & Discussion",
         Inches(3.5), Inches(4.35), Inches(6.3), Inches(0.6),
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s8, "We're happy to dive into any module,\ndata methodology, or design decision.",
         Inches(3.5), Inches(4.95), Inches(6.3), Inches(1.2),
         size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/Celia/Desktop/courses/quarter_2/pra_before/practicum_demo/Flex-Practicum-Project-2026/PPT/EMS_Intelligence_Platform_Showcase.pptx"
prs.save(out)
print(f"Saved: {out}")
